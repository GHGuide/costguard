#!/usr/bin/env python3
"""Fill the UiPath AgentHack template with CostGuard content — in place, no redesign.

  python3 deck/build_deck.py   # deck/template.pptx -> deck/costguard.pptx

Replaces every sample/Lorem-Ipsum string, fills the benefits table, drops in the
logo + dark architecture + hero-chart images, and keeps every layout / master /
brand element. Numbers come from BRIEF.md and the verified live runs.
"""
import os
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

HERE = os.path.dirname(os.path.abspath(__file__))
A = lambda f: os.path.join(HERE, "assets", f)
GREY = RGBColor(0x8A, 0x8A, 0x8A)

prs = Presentation(os.path.join(HERE, "template.pptx"))
S = list(prs.slides)


def shape(slide, sid):
    for sh in slide.shapes:
        if sh.shape_id == sid:
            return sh
    raise KeyError(sid)


def set_lines(sh, lines, tag=None):
    """Replace a text frame's content paragraph-by-paragraph, preserving the first
    run's formatting per paragraph. Optional small grey criterion tag as a final line."""
    tf = sh.text_frame
    all_lines = list(lines) + ([tag] if tag else [])
    while len(tf.paragraphs) < len(all_lines):
        tf._txBody.append(deepcopy(tf.paragraphs[-1]._p))
    for i, line in enumerate(all_lines):
        p = tf.paragraphs[i]
        runs = p.runs
        if runs:
            runs[0].text = line
            for extra in runs[1:]:
                extra._r.getparent().remove(extra._r)
        else:
            p.add_run().text = line
    for p in tf.paragraphs[len(all_lines):]:
        p._p.getparent().remove(p._p)
    if tag:
        r = tf.paragraphs[len(all_lines) - 1].runs[0]
        r.font.size = Pt(10); r.font.italic = True
        try: r.font.color.rgb = GREY
        except Exception: pass


def set_cell(tbl, r, text):
    tf = tbl.cell(r, 1).text_frame
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for extra in p.runs[1:]:
            extra._r.getparent().remove(extra._r)
    else:
        p.add_run().text = text
    for pp in tf.paragraphs[1:]:
        pp._p.getparent().remove(pp._p)


def drop_picture(slide, path, left, top, width):
    slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))


# ---- Slide 0: cover (keep the AgentHack line #168; replace #169) ----
set_lines(shape(S[0], 169), [
    "CostGuard",
    "Cost-regression testing for AI agents",
])

# ---- Slide 1: was the sample team's baked photo slide -> solo builder card ----
pic = shape(S[1], 176)
pic._element.getparent().remove(pic._element)          # remove baked "Jane Doe" image
set_lines(shape(S[1], 175), ["CostGuard"])
drop_picture(S[1], A("cg-logo.png"), 5.17, 1.95, 3.0)
tb = S[1].shapes.add_textbox(Inches(3.5), Inches(5.25), Inches(6.33), Inches(0.7))
tb.text_frame.word_wrap = True
p = tb.text_frame.paragraphs[0]
from pptx.enum.text import PP_ALIGN
p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Built solo by Leonardo Calancea  ·  Track 3 (Test Cloud)"
r.font.size = Pt(16); r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

# ---- Slide 2: Problem | Solution ----
set_lines(shape(S[2], 184), [
    "Teams ship AI agents fast. A small prompt or model change can quietly triple token spend while still passing every correctness test.",
    "QA checks whether the agent is right, never what it costs to be right. The bills are board-level: FinOps X 2026 called it the Great Token Panic, and one 4-agent loop burned about $47K in 11 days.",
    "There is a correctness gate before production. There is no cost gate.",
], tag="Scores: Business Impact")
set_lines(shape(S[2], 186), [
    "CostGuard adds a first-class test type: cost-regression testing, run and governed on UiPath.",
    "Before a change ships, it runs the agent many times on a Test Cloud scenario set, measures cost per correctly-processed invoice (not per token), compares to the last approved baseline, and blocks a regression. The call goes to a human in Action Center.",
    "Cost is paired with a quality score, so a cheaper-but-worse agent cannot slip through. The line nobody else can draw: cost per business outcome.",
], tag="Scores: Creativity · Platform Usage")

# ---- Slide 3: Benefits + technologies table ----
set_lines(shape(S[3], 196), [
    "Blocks silent cost regressions before production. In the demo a “smarter” candidate scored +4.5% accuracy but cost 7× more per processed invoice, and CostGuard blocked it. On real UiPath-gateway models the same trap cost 13× for no accuracy gain.",
    "It prices agents in dollars per business outcome, works on any agent or framework, and a human always owns the final call.",
], tag="Scores: Business Impact · Technical Execution")
tbl = shape(S[3], 193).table
set_cell(tbl, 0, "Platform, QA and FinOps engineers shipping AI agents")
set_cell(tbl, 1, "Engineering, QA, Platform Ops, FinOps")
set_cell(tbl, 2, "Enterprises running production AI agents (finance, insurance, BPO)")
set_cell(tbl, 3, "Test Cloud, Coded Agents (Python SDK), Orchestrator, Action Center, Document Understanding, AI Trust Layer Gateway")
set_cell(tbl, 4, "LangChain / CrewAI agent-under-test, OpenTelemetry, Claude Code via uip skills")

# ---- Slide 4: architecture (keep title #204; caption #205 + image) ----
set_lines(shape(S[4], 205), [
    "A change goes in. CostGuard runs the agent many times on UiPath, prices each correct outcome, compares to the baseline, and blocks a regression. A human owns any block.",
], tag="Scores: Platform Usage · Technical Execution")
drop_picture(S[4], A("cg-arch.png"), 0.7, 3.35, 11.9)

# ---- Slide 5: hero / live proof ----
set_lines(shape(S[5], 212), ["Live proof: real, not a slide"])
set_lines(shape(S[5], 213), [
    "A live serverless job on UiPath Automation Cloud returned verdict FAIL: cost 7.26× for +4.5% accuracy. On the AI Trust Layer gateway the same upgrade cost 13.12× for zero gain.",
    "The verdict is registered as a first-class Test Cloud result, and a 30-scenario regression suite keeps the gate's own verdicts pinned in CI.",
], tag="Scores: Completeness · Technical Execution")
drop_picture(S[5], A("cg-chart.png"), 3.42, 3.3, 6.5)

# ---- Slide 6: closing ----
set_lines(shape(S[6], 218), [
    "Looked fine. Cost 7× more. Blocked before it shipped.",
    "CostGuard  ·  github.com/GHGuide/costguard",
])

out = os.path.join(HERE, "CG-submission.pptx")
prs.save(out)
print("saved", out, "|", len(prs.slides), "slides")
